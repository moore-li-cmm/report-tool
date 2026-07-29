#!/usr/bin/env python3
"""Render the exec summary as a native, editable single-slide .pptx.

Takes data.json + narrative.json and produces a slide meant to be shared and
hand-tweaked in PowerPoint/Keynote/Google Slides before sending on. Everything
is real text boxes / shapes / a native chart, so it edits like any normal
slide. Self-contained: only python-pptx, no Jira engine or HTML template.

Layout follows the Innova "Executive Summary" format: a KPI tile row, then
Active epics + What's-next (left), Key updates (center), and Throughput trend +
Focus areas (right). Scoped to the fields PART's Jira data actually supports
(no Spend / Say-Do — Story points/Velocity render once a sprint exists, else a
"not tracked"/"awaiting data" placeholder, never a mock number).

Usage:
    ./.venv/bin/python .claude/skills/data-format-report/scripts/format_pptx.py
"""

from __future__ import annotations

import json
import os
from datetime import datetime

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --- palette (approximating the Innova reference) --------------------------
HEADER_GREEN = RGBColor(0x1B, 0x5E, 0x4F)   # section header bars
TILE_BORDER = RGBColor(0x53, 0x9E, 0x6B)    # KPI tile outline
VALUE_GREEN = RGBColor(0x1B, 0x5E, 0x4F)    # KPI value text
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_GRAY = RGBColor(0x59, 0x57, 0x53)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BAR_GOOD = RGBColor(0x2E, 0x8B, 0x3D)
BAR_MID = RGBColor(0xE0, 0x9B, 0x00)
BAR_LOW = RGBColor(0x2A, 0x78, 0xD6)
# Epic priority swatch, keyed by rank (5 = Highest … 1 = Lowest, 0 = none/unknown).
PRIORITY_COLORS = {
    5: RGBColor(0xC0, 0x39, 0x39),
    4: RGBColor(0xE0, 0x9B, 0x00),
    3: RGBColor(0x2A, 0x78, 0xD6),
    2: RGBColor(0x53, 0x9E, 0x6B),
    1: RGBColor(0x8A, 0x88, 0x82),
    0: RGBColor(0xB0, 0xB0, 0xAA),
}

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.3)

# Repo root is four levels up (scripts/ -> data-format-report/ -> skills/ -> .claude/).
# The pipeline reads data.json + narrative.json and writes exec_summary.pptx there.
_REPO_ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "..", ".."))
DATA_PATH = os.path.join(_REPO_ROOT, "data.json")
NARRATIVE_PATH = os.path.join(_REPO_ROOT, "narrative.json")
OUT_PATH = os.path.join(_REPO_ROOT, "exec_summary.pptx")


def _tf(shape):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    return tf


def _para(tf, text, size, *, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT,
          first=False, space_after=2):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p


def _run(p, text, size, *, bold=False, color=TEXT_DARK, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return r


def section_header(slide, x, y, w, text):
    h = Inches(0.3)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = HEADER_GREEN
    box.line.fill.background()
    box.shadow.inherit = False
    tf = _tf(box)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _para(tf, text.upper(), 10.5, bold=True, color=WHITE, first=True)
    return Emu(y + h)


def kpi_tile(slide, x, y, w, h, value, label, *, empty_dot=False, sub=None, value_size=19):
    tile = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    tile.fill.solid()
    tile.fill.fore_color.rgb = WHITE
    tile.line.color.rgb = TILE_BORDER
    tile.line.width = Pt(1.5)
    tile.shadow.inherit = False
    tf = _tf(tile)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    if empty_dot:
        # Deliberately blank: an unfilled circle for whoever presents the slide
        # to color/annotate by hand in PowerPoint. Never auto-computed — this
        # tile is a manual fill-in, not a data-driven one.
        d = Inches(0.32)
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Emu(x + (w - d) // 2), Emu(y + Inches(0.16)), d, d
        )
        dot.fill.background()
        dot.line.color.rgb = TILE_BORDER
        dot.line.width = Pt(1.5)
        dot.shadow.inherit = False
        _para(tf, "", 6, first=True)
        _para(tf, label.upper(), 8, bold=True, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
    else:
        _para(tf, value, value_size, bold=True, color=VALUE_GREEN, align=PP_ALIGN.CENTER, first=True)
        _para(tf, label.upper(), 8, bold=True, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
        if sub:
            _para(tf, sub, 7.5, color=TEXT_GRAY, align=PP_ALIGN.CENTER)


def epic_row(slide, x, y, w, epic):
    """One active-epic line: priority swatch + label, the epic name with NEW /
    priority-change badges, and a right-aligned status (in-flight emphasized).
    Epics are pre-sorted by priority upstream, so top-to-bottom IS priority order."""
    h = Inches(0.30)
    rank = epic.get("priority_rank", 0)
    pri_color = PRIORITY_COLORS.get(rank, PRIORITY_COLORS[0])
    dot_w = Inches(0.12)
    right_w = Inches(1.2)

    # priority swatch (color = rank; the text label carries the exact level)
    dot = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Emu(y + Inches(0.05)), dot_w, Inches(0.2))
    dot.fill.solid()
    dot.fill.fore_color.rgb = pri_color
    dot.line.fill.background()
    dot.shadow.inherit = False

    # name + inline badges
    name_x = Emu(x + dot_w + Inches(0.06))
    name_w = Emu(w - dot_w - Inches(0.06) - right_w)
    nb = slide.shapes.add_textbox(name_x, y, name_w, h)
    tf = _tf(nb)
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.space_after = Pt(0)
    _run(p, f"{epic.get('priority') or '—'}  ", 7, bold=True, color=pri_color)
    name = f"{epic.get('key', '')}: {epic.get('summary', '')}"
    _run(p, (name[:30] + "…") if len(name) > 31 else name, 8, color=TEXT_DARK)
    if epic.get("is_new"):
        _run(p, "  NEW", 7, bold=True, color=BAR_GOOD)
    pc = epic.get("priority_change")
    if pc:
        arrow = {"raised": "▲", "lowered": "▼"}.get(pc.get("direction"), "⇄")
        _run(p, f"  {arrow}PRI", 7, bold=True, color=BAR_MID)

    # right: current status + child progress, in-flight emphasized
    rb = slide.shapes.add_textbox(Emu(x + w - right_w), y, right_w, h)
    tf = _tf(rb)
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p.space_after = Pt(0)
    in_flight = epic.get("in_flight")
    total, done = epic.get("total", 0), epic.get("done", 0)
    prog = f"  {done}/{total}" if total else ""
    # Normalize the in-flight state to a plain-language label instead of the raw
    # Jira status name (e.g. "Analyzing" -> "In Progress").
    label = "In Progress" if in_flight else epic.get("status", "")
    _run(p, f"{label}{prog}", 7.5, bold=bool(in_flight),
         color=HEADER_GREEN if in_flight else TEXT_GRAY)


def bullets(slide, x, y, w, h, items, size=9.5):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = _tf(box)
    # Shrink text to fit the fixed box so a wordy panel can't overflow downward
    # behind the panel stacked below it (PowerPoint recalcs the fit on open).
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    if not items:
        _para(tf, "—", size, color=TEXT_GRAY, first=True)
        return
    for i, item in enumerate(items):
        _para(tf, f"•  {item}", size, color=TEXT_DARK, first=(i == 0), space_after=5)


def trend_chart(slide, x, y, w, h, trend, annotation):
    counts = trend.get("counts", [])
    labels = trend.get("week_labels", [])
    cats = []
    for lab in labels:
        try:
            cats.append(datetime.strptime(lab, "%Y-%m-%d").strftime("%-m/%-d"))
        except ValueError:
            cats.append(lab)
    cd = CategoryChartData()
    cd.categories = cats or ["—"]
    cd.add_series("Resolved / wk", counts or [0])
    gf = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x, y, w, h, cd)
    chart = gf.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    plot.series[0].format.line.color.rgb = BAR_LOW
    for axis in (chart.category_axis, chart.value_axis):
        axis.tick_labels.font.size = Pt(7)
        axis.tick_labels.font.color.rgb = TEXT_GRAY
    if annotation:
        note = slide.shapes.add_textbox(x, Emu(y + h), w, Inches(0.4))
        tf = _tf(note)
        _para(tf, annotation, 7.5, color=TEXT_GRAY, first=True)


def velocity_chart(slide, x, y, w, h, velocity):
    """Committed vs. completed story points per closed sprint. Only called once
    velocity_history is non-null (at least one closed sprint) — see the
    "Awaiting data" fallback in build() for the empty case."""
    cd = CategoryChartData()
    cd.categories = velocity.get("sprint_labels", [])
    cd.add_series("Committed", velocity.get("committed", []))
    cd.add_series("Completed", velocity.get("completed", []))
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, cd)
    chart = gf.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(7)
    chart.has_title = False
    plot = chart.plots[0]
    plot.series[0].format.fill.solid()
    plot.series[0].format.fill.fore_color.rgb = TILE_BORDER
    plot.series[1].format.fill.solid()
    plot.series[1].format.fill.fore_color.rgb = BAR_LOW
    for axis in (chart.category_axis, chart.value_axis):
        axis.tick_labels.font.size = Pt(7)
        axis.tick_labels.font.color.rgb = TEXT_GRAY


def build(data: dict, narrative: dict, out_path: str) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # --- title + status line ---
    title = slide.shapes.add_textbox(MARGIN, Inches(0.12), Inches(9), Inches(0.55))
    tf = _tf(title)
    _para(tf, f"Executive Summary — {data.get('project','')}", 26, bold=True,
          color=HEADER_GREEN, first=True)
    sub = slide.shapes.add_textbox(MARGIN, Inches(0.66), Inches(12.7), Inches(0.3))
    tf = _tf(sub)
    _para(tf, f"{data.get('period_label','')}  ·  generated {data.get('generated_at','')}",
          9.5, color=TEXT_GRAY, first=True)

    # --- KPI tiles ---
    since_days = data.get("since_days", 30)
    epics = data.get("recent_epics", {}).get("linked", [])
    epics_started = sum(1 for e in epics if e.get("is_new"))
    epics_done = sum(1 for e in epics if e.get("is_done_recent"))
    stories_completed = sum(
        1 for i in data.get("resolved_this_period", []) if i.get("issuetype", "").lower() == "story"
    )

    ec = data.get("epic_cycle_time") or {}
    ecd = ec.get("days")
    tiles = [
        # Project health: deliberately blank — a manual fill-in the presenter
        # colors/annotates by hand in PowerPoint, never auto-computed.
        {"label": "Project health", "empty_dot": True},
        {"value": f"{epics_started} | {epics_done}", "label": "Started epics | done",
         "sub": f"Last {since_days}d"},
        {"value": f"{data.get('backlog_total',0)} | {data.get('backlog_delivered',0)}", "label": "Backlog | delivered"},
        {"value": f"{ecd if ecd is not None else 'n/a'} d", "label": "Epic cycle time"},
        {"value": str(stories_completed), "label": "Stories completed"},
        # sprint_goal is null until a sprint is active (see jira_exec_summary
        # .compute_sprint_stats) — show real committed/completed points once one
        # is, else the same "not tracked" placeholder (matches the velocity
        # panel below). No mock number either way.
        (
            {"value": f"{int(sg['completed_points'])} | {int(sg.get('in_progress_points', 0))} | {int(sg['committed_points'])}",
             "label": "Story pts: done | WIP | total", "sub": sg["name"], "value_size": 15}
            if (sg := data.get("sprint_goal"))
            else {"value": "—", "label": "Story points", "sub": "no active sprint"}
        ),
        {"value": str(len(data.get("blockers", []))), "label": "Flagged"},
    ]
    pr = data.get("pull_requests") or {}
    if pr.get("configured"):
        # Show "merged this period | open now" — the two numbers a reader can
        # reconcile against the repo (open_now matches GitHub's open count).
        # Deliberately NOT "opened", which collides with "open" (created-in-window
        # vs currently-open are different things and confuse readers).
        open_now = pr.get("open_now", 0)
        cur = pr.get("current_sprint")
        bucket = (pr.get("by_sprint") or {}).get(cur) if cur else None
        pr_tile = ({"value": f"{bucket['merged']} | {open_now}", "label": f"PRs merged ({cur}) | open now"}
                   if bucket else
                   {"value": f"{pr.get('merged_in_window', 0)} | {open_now}", "label": f"PRs merged {since_days}d | open now"})
        tiles.insert(len(tiles) - 1, pr_tile)  # just before Flagged
    n = len(tiles)
    gap = Inches(0.12)
    usable = SLIDE_W - 2 * MARGIN - (n - 1) * gap
    tw = Emu(usable // n)
    ty = Inches(1.05)
    th = Inches(0.92)
    for i, t in enumerate(tiles):
        tx = Emu(MARGIN + i * (tw + gap))
        kpi_tile(slide, tx, ty, tw, th, t.get("value"), t["label"],
                 empty_dot=t.get("empty_dot", False), sub=t.get("sub"), value_size=t.get("value_size", 19))

    # --- three body columns; each has a top slot and a bottom slot ---
    body_y = Inches(2.2)
    bottom_y = Inches(5.05)  # pushed down so a wordy Key-updates block clears Focus areas
    lx, lw = MARGIN, Inches(4.0)
    cx, cw = Inches(4.45), Inches(4.35)
    rx, rw = Inches(8.95), Inches(4.08)

    # LEFT: active epics (top) + what's next (bottom). Rows are in priority order
    # (sorted upstream); each shows in-flight status and NEW / priority-change badges.
    y = section_header(slide, lx, body_y, lw, "Active epics — by priority")
    ey = Emu(y + Inches(0.06))
    for e in epics[:6]:
        epic_row(slide, lx, ey, lw, e)
        ey = Emu(ey + Inches(0.30))
    if narrative.get("whats_next"):
        wy = section_header(slide, lx, bottom_y, lw, "What's next")
        bullets(slide, lx, Emu(wy + Inches(0.06)), lw, Inches(2.0),
                narrative["whats_next"], size=9)

    # CENTER: key updates (top) + focus areas (directly below it)
    y = section_header(slide, cx, body_y, cw, "Key updates")
    bullets(slide, cx, Emu(y + Inches(0.06)), cw, Inches(2.4), narrative.get("key_updates", []), size=9)
    fy = section_header(slide, cx, bottom_y, cw, "Focus areas")
    bullets(slide, cx, Emu(fy + Inches(0.06)), cw, Inches(2.0), narrative.get("focus_areas", []), size=9)

    # RIGHT: throughput trend (top) + velocity (bottom). Velocity renders once
    # velocity_history is non-null (jira_exec_summary.compute_sprint_stats) —
    # i.e. at least one sprint has closed — else the "awaiting data" fallback.
    y = section_header(slide, rx, body_y, rw, "Throughput trend")
    trend_chart(slide, rx, Emu(y + Inches(0.06)), rw, Inches(1.9),
                data.get("trend", {}), narrative.get("trend_annotation"))
    vy = section_header(slide, rx, bottom_y, rw, "Velocity — story points / sprint")
    velocity = data.get("velocity_history")
    if velocity:
        velocity_chart(slide, rx, Emu(vy + Inches(0.06)), rw, Inches(1.9), velocity)
    else:
        ph = slide.shapes.add_textbox(rx, Emu(vy + Inches(0.06)), rw, Inches(2.0))
        tf = _tf(ph)
        p = _para(tf, "Awaiting data — needs at least one completed sprint with "
                      "estimated story points. No placeholder numbers shown until "
                      "one closes.",
                  9, color=TEXT_GRAY, first=True)
        p.runs[0].font.italic = True

    # (Data-note footnote intentionally omitted from the slide — caveats still
    # live in data.json/auto_caveats for reference.)

    prs.save(out_path)


def main() -> None:
    with open(DATA_PATH, encoding="utf-8") as f:
        data = json.load(f)
    with open(NARRATIVE_PATH, encoding="utf-8") as f:
        narrative = json.load(f)

    build(data, narrative, OUT_PATH)
    print(f"Wrote {OUT_PATH}")


if __name__ == "__main__":
    main()

#!/usr/bin/env python3
"""Render a reference slide showing HOW each field on exec_summary.pptx is
derived — same layout/geometry as format_pptx.py, but every tile/panel holds
an explanation (source query, function, file) instead of real data.

Not part of the reporting pipeline: nothing reads data.json/narrative.json
here, and no other script depends on this one's output. Exists purely as a
documentation artifact — a live audit of where each number on the real slide
comes from, in the same spatial layout so it's easy to compare side by side.

Usage:
    ./.venv/bin/python .claude/skills/data-format-report/scripts/explain_layout.py
"""

from __future__ import annotations

import os

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Emu, Inches

from format_pptx import (
    BODY_Y,
    BOTTOM_Y,
    HEADER_GREEN,
    MARGIN,
    SLIDE_H,
    SLIDE_W,
    TEXT_GRAY,
    TILE_ROW_H,
    TILE_ROW_Y,
    _para,
    _REPO_ROOT,
    _tf,
    body_columns,
    bullets,
    section_header,
    tile_card,
    tile_row_geometry,
)

OUT_PATH = os.path.join(_REPO_ROOT, "exec_summary_explainer.pptx")

TILES = [
    ("Project health",
     "Deliberately blank -- an empty, uncolored circle. Never auto-computed and never "
     "written by the manager into narrative.json; a manual fill-in for whoever presents "
     "the slide to annotate by hand in PowerPoint. suggested_status still exists in "
     "data.json as a naive reference signal (flagged->critical, nothing delivered->warning, "
     "else good), but nothing here reads it."),
    ("Started epics | done",
     "is_new count | is_done_recent count, both over recent_epics.linked (compute_epics). "
     "is_new = epic created within the 30-day window; is_done_recent = epic resolved within "
     "the 30-day window. Both halves are genuinely time-boxed to the same window (sub-label)."),
    ("Backlog | delivered",
     "backlog_total = OPEN tickets whose parent epic rolls up to a real Initiative "
     "(snapshot, NOT sprint-scoped -- differs from Jira's own board Backlog panel, which "
     "excludes sprinted tickets instead). backlog_delivered = tickets RESOLVED in the last "
     "30 days AND whose parent epic rolls up to a real Initiative -- same scope as "
     "backlog_total, so the two numbers reconcile. Orphan/out-of-scope tickets on either "
     "side are excluded and counted in auto_caveats; can legitimately read 0 delivered even "
     "when real work shipped outside the initiative's tracked epics."),
    ("Epic cycle time",
     "Days from EPIC creation -> resolution, averaged over epics resolved in the current "
     "30-day window (compute_epic_cycle_time). Discard-status epics excluded; null when 0 "
     "epics resolved. No sub-value/delta shown on the tile (prior_days is in data.json for "
     "the manager to use in prose only)."),
    ("Stories completed",
     "Count of data.resolved_this_period items where issuetype == 'Story' (that list is "
     "already scoped to initiative-epic children -- see Backlog | delivered). A raw 30-day "
     "count, not a per-week rate."),
    ("Story pts: done | WIP | total",
     "Sum of Story Points (customfield_13078) on tickets carrying the active Sprint "
     "(customfield_10020). done = resolved tickets' points; WIP = unresolved tickets whose "
     "Jira status category is 'In Progress'; total = all points committed to the sprint. "
     "Shows '-' until a sprint is active (compute_sprint_stats)."),
    ("PRs merged (sprint) | open now",
     "GitHub PR counts via GITHUB_TOKEN/GITHUB_REPOS (github_prs.py). merged = PRs merged "
     "during the linked ticket's sprint (or last 30 days if no sprint yet); open_now = the "
     "repo's live open-PR count. PRs cross-link to Jira keys parsed from PR title/body."),
    ("Flagged",
     "Count of issues carrying Jira's native Flagged/Impediment field (customfield_10021 on "
     "this instance) whose own status is not Done-like (Done/Closed/Resolved/Discard/"
     "Cancelled). Same signal as the flag icon on the board. Scanned project-wide, not just "
     "initiative-linked tickets."),
]

LEFT_TOP = (
    "ACTIVE EPICS -- BY RANK",
    [
        "Source: recent_epics.linked -- epics whose Jira parent = Initiative AA-431.",
        "Sort: Jira's native Rank field (customfield_10019, LexoRank), ascending, "
        "engine-sorted -- list order IS the team's real backlog order. Epics carry no "
        "priority field at all: Priority sits at an unused default (\"Lowest\") on every "
        "PART epic and carries no signal.",
        "Each row: ordinal #N rank badge (position in the sorted list); NEW badge = created "
        "within the 30-day window; up/down-RANK badge = Rank moved within the window (via "
        "Jira changelog -- direction only, Jira logs no absolute from/to position); right "
        "side = status (Jira 'In Progress' category shown as 'In Progress') + done/total "
        "child tickets.",
    ],
)
LEFT_BOTTOM = (
    "WHAT'S NEXT",
    [
        "Manually written by the manager subagent (narrative.json -> whats_next) -- NOT auto-computed.",
        "Derived from epic sequencing/descriptions and pending decisions in data.json; "
        "honest timing only ('next sprint'), never a fabricated date.",
    ],
)
RIGHT_TOP = (
    "KEY UPDATES",
    [
        "Manually written (narrative.json -> key_updates).",
        "Must be grounded in data.resolved_this_period -- tickets RESOLVED in the 30-day "
        "window, using real description text (ADF flattened via jira_report.adf_to_text).",
        "Never based on epic goals/aspirational scope, and never inflated past what the "
        "ticket supports.",
    ],
)
RIGHT_BOTTOM = (
    "FOCUS AREAS",
    [
        "Manually written (narrative.json -> focus_areas).",
        "Drawn from data.flagged (Jira's native Flag/Impediment field), data.stale (open >=14 days untouched), "
        "data.overdue (past due date) -- surfaces the worst by days_since_update/days_overdue.",
        "No dedicated trend-chart panel exists -- if there's a real improvement/decline "
        "story in data.trend, fold it in here or in key_updates instead.",
    ],
)


def explain_tile(slide, x, y, w, h, label, explanation):
    tf = tile_card(slide, x, y, w, h)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    _para(tf, label.upper(), 7.5, bold=True, color=HEADER_GREEN, first=True, space_after=1)
    _para(tf, explanation, 6.2, color=TEXT_GRAY, space_after=0)


def panel(slide, x, y, w, h, title, lines):
    hy = section_header(slide, x, y, w, title)
    bullets(slide, x, Emu(hy + Inches(0.06)), w, h, lines, size=8.5)


def build(out_path: str) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    title = slide.shapes.add_textbox(MARGIN, Inches(0.12), Inches(11), Inches(0.55))
    tf = _tf(title)
    _para(tf, "Executive Summary -- HOW EACH FIELD IS CALCULATED", 22, bold=True,
          color=HEADER_GREEN, first=True)
    sub = slide.shapes.add_textbox(MARGIN, Inches(0.62), Inches(12.7), Inches(0.35))
    tf = _tf(sub)
    _para(tf, "Reference copy of exec_summary.pptx's layout. Not real data -- each "
              "element explains its source query/computation instead. Regenerate with "
              "explain_layout.py whenever the pipeline's field definitions change.",
          9, color=TEXT_GRAY, first=True)

    # --- KPI tile row (shared geometry from format_pptx) ---
    n = len(TILES)
    tw, gap = tile_row_geometry(n)
    ty = TILE_ROW_Y
    th = TILE_ROW_H
    for i, (label, explanation) in enumerate(TILES):
        tx = Emu(MARGIN + i * (tw + gap))
        explain_tile(slide, tx, ty, tw, th, label, explanation)

    # --- two body columns (shared geometry from format_pptx); each has a top
    # slot and a bottom slot, split evenly across the slide width ---
    body_y = BODY_Y
    bottom_y = BOTTOM_Y
    lx, rx, col_w = body_columns()
    lw = rw = col_w

    panel(slide, lx, body_y, lw, Inches(2.6), *LEFT_TOP)
    panel(slide, lx, bottom_y, lw, Inches(2.0), *LEFT_BOTTOM)
    panel(slide, rx, body_y, rw, Inches(2.6), *RIGHT_TOP)
    panel(slide, rx, bottom_y, rw, Inches(2.0), *RIGHT_BOTTOM)

    prs.save(out_path)


def main() -> None:
    build(OUT_PATH)
    print(f"Wrote {OUT_PATH}")


if __name__ == "__main__":
    main()
